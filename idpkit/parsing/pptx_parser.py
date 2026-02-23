"""PowerPoint parser using python-pptx for text extraction."""

import logging

from .base import BaseParser, ParseResult

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """Extract text from PowerPoint presentations.

    Uses ``python-pptx`` to iterate over slides, extracting text from shapes
    and slide notes.  Each slide becomes a separate "page" in the
    :class:`ParseResult`.
    """

    def supported_extensions(self) -> list[str]:
        return ["pptx", "ppt"]

    def parse(self, file_path: str) -> ParseResult:
        """Parse a PPTX file and return its text content.

        Parameters
        ----------
        file_path:
            Path to the ``.pptx`` file.

        Raises
        ------
        ImportError
            If python-pptx is not installed.
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX parsing. "
                "Install it with: pip install python-pptx"
            )

        logger.info("Parsing PPTX: %s", file_path)

        prs = Presentation(file_path)
        pages: list[dict] = []
        all_text_parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            slide_parts.append(f"## Slide {slide_num}")

            # Extract text from all shapes on the slide.
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)

                # Handle tables in slides.
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip().replace("|", "").strip():
                            slide_parts.append(row_text)

            # Extract slide notes.
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n**Notes:** {notes_text}")

            slide_text = "\n".join(slide_parts)
            pages.append({"page": slide_num, "text": slide_text})
            all_text_parts.append(slide_text)

        page_count = len(pages)
        full_text = "\n\n".join(all_text_parts)

        # Metadata
        metadata: dict = {}
        core = prs.core_properties
        if core.title:
            metadata["title"] = core.title
        if core.author:
            metadata["author"] = core.author
        if core.subject:
            metadata["subject"] = core.subject
        metadata = {k: v for k, v in metadata.items() if v}
        metadata["slide_count"] = page_count

        logger.info("PPTX parsing complete: %d slides extracted", page_count)

        return ParseResult(
            text=full_text,
            pages=pages,
            metadata=metadata,
            page_count=page_count,
        )
